from pptx import Presentation

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Bullet slide layout'

p = tf.add_paragraph()
p.text = 'First bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Subsequent bullet'
p.level = 2

p = tf.add_paragraph()
p.text = 'Another bullet'
p.level = 2

prs.save('bullet.pptx')