
import sys, getopt
import math
from pptx import Presentation
from pptx.util import Inches

def read_two_per_line(inputfile, outputfile, delim=';'):
	prs = Presentation()
	title_only_slide_layout = prs.slide_layouts[5]

	lines = [line.strip() for line in open(inputfile)]
	print("Read %d lines from %s" % (len(lines), inputfile))
	nfile = len(lines)

	nlines = 5
	npages = nslides = int(len(lines)/(nlines))
	nblock = nslides*nlines
	rlines = len(lines) - nblock
	if rlines > 0:
		npages += 1
	# print("First %d lines" % nblock)
	for n in range(0, nslides):
		print("Slide %d of %d" % (n+1, npages))
		slide = prs.slides.add_slide(title_only_slide_layout)
		shapes = slide.shapes

		shapes.title.text = 'Page %d of %d' % (n+1, npages)

		rows = nlines+1
		cols = 2
		left = top = Inches(2.0)
		width = Inches(6.0)
		height = Inches(4.8)

		table = shapes.add_table(rows, cols, left, top, width, height).table

		# set column widths
		table.columns[0].width = Inches(2.0)
		table.columns[1].width = Inches(4.0)

		# write column headings
		table.cell(0, 0).text = "One"
		table.cell(0, 1).text = "Two"

		for i in range(0, nlines):
			# write body cells
			x = n*nlines+i
			# print("Line %d : record %d of %d" %(i, x, nfile))
			names = lines[x].split(delim)
			j=0
			for name in names:
				print("\t%30s" % (name), end='')
				table.cell(i+1, j).text = name
				j += 1
			print('')

	if rlines > 0:
		# print("%d lines left" % rlines)	
		print("Slide %d of %d" % (npages, npages))
		slide = prs.slides.add_slide(title_only_slide_layout)
		shapes = slide.shapes

		shapes.title.text = 'Page %d of %d' % (npages, npages)

		rows = math.ceil(rlines)+1
		cols = 2
		left = top = Inches(2.0)
		width = Inches(6.0)
		height = Inches(4.8)

		table = shapes.add_table(rows, cols, left, top, width, height).table

		# set column widths
		table.columns[0].width = Inches(2.0)
		table.columns[1].width = Inches(4.0)

		# write column headings
		table.cell(0, 0).text = "One"
		table.cell(0, 1).text = "Two"
		
		for i in range(0, nfile-nblock):
			# write body cells
			# print("Extra line %d of %d : " %(i+1, nfile-nblock))
			names = lines[nblock+i].split(delim)
			j = 0
			for name in names:
				print("\t%30s" % (name), end='')
				table.cell(i+1, j).text = name
				j += 1
			print('')
			
	print("Write %s" % outputfile)
	prs.save(outputfile)


def read_one_per_line(inputfile, outputfile):
	prs = Presentation()
	title_only_slide_layout = prs.slide_layouts[5]

	lines = [line.strip() for line in open(inputfile)]
	print("Read %d lines from %s" % (len(lines), inputfile))
	nfile = len(lines)

	nlines = 5
	npages = nslides = int(len(lines)/(nlines*2))
	nblock = nslides*nlines*2
	rlines = len(lines) - nblock
	if rlines > 0:
		npages += 1
		
	# print("First %d lines" % nblock)
	for n in range(0, nslides):
		print("Slide %d of %d" % (n+1, npages))
		slide = prs.slides.add_slide(title_only_slide_layout)
		shapes = slide.shapes

		shapes.title.text = 'Page %d of %d' % (n+1, npages)

		rows = nlines+1
		cols = 2
		left = top = Inches(2.0)
		width = Inches(6.0)
		height = Inches(4.8)

		table = shapes.add_table(rows, cols, left, top, width, height).table

		# set column widths
		table.columns[0].width = Inches(2.0)
		table.columns[1].width = Inches(4.0)

		# write column headings
		table.cell(0, 0).text = "One"
		table.cell(0, 1).text = "Two"

		for i in range(0, nlines):
			# write body cells
			x = n*nlines*2+i*2
			y = n*nlines*2+i*2+1
			# print("Line %d : record %d and %d of %d" %(i, x, y, nfile))
			print("\t%30s : %s" % (lines[x], lines[y]))
			table.cell(i+1, 0).text = lines[x]
			table.cell(i+1, 1).text = lines[y]

	if rlines > 0:
		# print("%d lines left" % rlines)	
		print("Slide %d of %d" % (npages, npages))
		slide = prs.slides.add_slide(title_only_slide_layout)
		shapes = slide.shapes

		shapes.title.text = 'Page %d of %d' % (npages, npages)

		rows = math.ceil(rlines/2)+1
		cols = 2
		left = top = Inches(2.0)
		width = Inches(6.0)
		height = Inches(4.8)

		table = shapes.add_table(rows, cols, left, top, width, height).table

		# set column widths
		table.columns[0].width = Inches(2.0)
		table.columns[1].width = Inches(4.0)

		# write column headings
		table.cell(0, 0).text = "One"
		table.cell(0, 1).text = "Two"
		j = 1
		for i in range(nslides*nlines*2, nfile, 2):
			# write body cells
			# print("Line %d and %d of %d" %(i, i+1, nfile))
			if i+1 >= nfile:
				print("\t%30s : %s" % (lines[i], "---"))
				table.cell(j, 0).text = lines[i]
				table.cell(j, 1).text = "---"
				break
			print("\t%30s : %s" % (lines[i], lines[i+1]))
			table.cell(j, 0).text = lines[i]
			table.cell(j, 1).text = lines[i+1]
			j += 1
	
	print("Write %s" % outputfile)
	prs.save(outputfile)
	

def main(argv):

	mode = 1
	inputfile = ''
	outputfile = ''
	try:
		opts, args = getopt.getopt(argv,"12hi:o:",["ifile=","ofile="])
	except getopt.GetoptError:
		print('test.py -i <inputfile> -o <outputfile>')
		sys.exit(2)
	for opt, arg in opts:
		if opt == '-h':
			print('test.py -i <inputfile> -o <outputfile>')
			sys.exit()
		if opt == '-1':
			mode = 1
		if opt == '-2':
			mode = 2
		elif opt in ("-i", "--ifile"):
			inputfile = arg
		elif opt in ("-o", "--ofile"):
			outputfile = arg
	if len(inputfile)==0 or len(outputfile)==0:
		print("Input and output files must be specified")
		return 1
		
	print('Input file is %s' % inputfile)
	print('Output file is %s' % outputfile)
	
	if mode == 2:
		read_two_per_line(inputfile, outputfile)
	else:
		read_one_per_line(inputfile, outputfile)
	
	return 0

	
if __name__ == "__main__":
	main(sys.argv[1:])
